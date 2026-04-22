from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_google_slide(prs, title_text, bullet_points, slide_type="text"):
    slide_layout = prs.slide_layouts[1] # Bullet layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Configurar Título
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(66, 133, 244) # Azul Google

    # Configurar Cuerpo
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.space_after = Pt(10)

def create_presentation():
    prs = Presentation()

    # --- SLIDE 1: PORTADA ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "LogiData: Decisiones Predictivas"
    subtitle.text = "Arquitectura de Datos de Alto Impacto\nIngeniero de Sistemas e Informática"
    
    # --- SLIDE 2: VISIÓN ---
    add_google_slide(prs, "Visión Estratégica", [
        "De Silos de Datos a Inteligencia de Negocio.",
        "Modelos predictivos de demanda basados en datos confiables.",
        "Validación automática de reglas de negocio en tiempo real."
    ])

    # --- SLIDE 3: ARQUITECTURA MEDALLION ---
    add_google_slide(prs, "Arquitectura Medallion", [
        "Bronze: Ingesta Raw (CSV) en S3.",
        "Silver: Calidad con Great Expectations y Delta Lake.",
        "Gold: Feature Store optimizado para Machine Learning."
    ])

    # --- SLIDE 4: CALIDAD ---
    add_google_slide(prs, "Blindaje de Calidad", [
        "Validación dinámica: Enums de zonas y rangos de precios.",
        "Gestión de fallos: Desvío automático a Cuarentena.",
        "Idempotencia: Operaciones Merge para evitar duplicados."
    ])

    # --- SLIDE 5: DELTA LAKE ---
    add_google_slide(prs, "El Poder de Delta Lake", [
        "Transacciones ACID: Integridad total en cada escritura.",
        "Time Travel: Auditoría histórica y reproducibilidad de modelos.",
        "Schema Enforcement: Protección contra datos mal estructurados."
    ])

    # --- SLIDE 6: AUTOMATIZACIÓN ---
    add_google_slide(prs, "Infraestructura como Código", [
        "Terraform: Despliegue replicable y versionado.",
        "Orquestación: Triggers de AWS Glue con dependencias lógicas.",
        "Escalabilidad: Serverless listo para crecer 10x."
    ])

    # --- SLIDE 7: FUTURO ML ---
    add_google_slide(prs, "Hacia el Machine Learning", [
        "Datasets maestros listos para entrenamiento.",
        "Analítica avanzada: Dashboards de alta velocidad.",
        "Enfoque en el sector bancario y logístico."
    ])

    prs.save('LogiData_Presentation.pptx')
    print("¡Presentación generada con éxito!")

create_presentation()